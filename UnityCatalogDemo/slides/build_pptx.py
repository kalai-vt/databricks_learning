"""Generate UnityCatalog_Demo_Slides.pptx from slide_content.md.

Usage: python3 build_pptx.py
Requires: pip install python-pptx
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE = Path(__file__).parent
SRC = HERE / "slide_content.md"
OUT = HERE / "UnityCatalog_Demo_Slides.pptx"

TITLE_COLOR = RGBColor(0xFF, 0x3B, 0x1D)   # Databricks-ish red-orange
TEXT_COLOR = RGBColor(0x1B, 0x1B, 0x1B)


def parse_slides(md_text: str):
    slides = []
    current_title = None
    current_bullets = []
    for line in md_text.splitlines():
        if line.startswith("## "):
            if current_title is not None:
                slides.append((current_title, current_bullets))
            current_title = line[3:].strip()
            current_bullets = []
        elif line.startswith("- "):
            current_bullets.append(line[2:].strip())
    if current_title is not None:
        slides.append((current_title, current_bullets))
    return slides


def build(slides):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for idx, (title, bullets) in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.1))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(36) if idx > 0 else Pt(44)
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR

        # Accent rule
        line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.35), Inches(3.0), Pt(4))
        line.fill.solid()
        line.fill.fore_color.rgb = TITLE_COLOR
        line.line.fill.background()

        # Bullets
        if bullets:
            body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.3))
            btf = body_box.text_frame
            btf.word_wrap = True
            for i, bullet in enumerate(bullets):
                bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                bp.text = f"•  {bullet}"
                bp.font.size = Pt(22)
                bp.font.color.rgb = TEXT_COLOR
                bp.space_after = Pt(12)

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(8), Inches(0.35))
        ftf = footer_box.text_frame
        fp = ftf.paragraphs[0]
        frun = fp.add_run()
        frun.text = "RetailCorp Unity Catalog Demo"
        frun.font.size = Pt(11)
        frun.font.color.rgb = RGBColor(0x8A, 0x8A, 0x8A)

        # Slide number
        num_box = slide.shapes.add_textbox(Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35))
        ntf = num_box.text_frame
        np_ = ntf.paragraphs[0]
        nrun = np_.add_run()
        nrun.text = str(idx + 1)
        nrun.font.size = Pt(11)
        nrun.font.color.rgb = RGBColor(0x8A, 0x8A, 0x8A)

    prs.save(OUT)
    print(f"Wrote {len(slides)} slides to {OUT}")


if __name__ == "__main__":
    slides = parse_slides(SRC.read_text())
    build(slides)
